"""Tests for elidia.tools.office — docx/xlsx/pptx read + write."""
from pathlib import Path

import pytest

from elidia.tools import ToolRegistry, create_default_registry
from elidia.tools.office import (
    _read_docx,
    _read_pptx,
    _read_xlsx,
    _write_docx,
    _write_xlsx,
    register_office_tools,
)


class TestRegistration:
    def test_registers_five_tools(self):
        registry = ToolRegistry()
        register_office_tools(registry)
        names = {t.name for t in registry.list_tools()}
        assert names == {"read_docx", "read_xlsx", "read_pptx", "write_docx", "write_xlsx"}

    def test_wired_into_default_registry(self):
        registry = create_default_registry()
        assert registry.get("read_docx") is not None


class TestDocxRoundTrip:
    @pytest.mark.asyncio
    async def test_write_then_read(self, tmp_dir: Path):
        path = tmp_dir / "report.docx"
        result = await _write_docx(str(path), "First paragraph.\nSecond paragraph.")
        assert not result.is_error
        assert path.exists()

        read_result = await _read_docx(str(path))
        assert not read_result.is_error
        assert "First paragraph." in read_result.content
        assert "Second paragraph." in read_result.content

    @pytest.mark.asyncio
    async def test_read_missing_file(self, tmp_dir: Path):
        result = await _read_docx(str(tmp_dir / "nope.docx"))
        assert result.is_error
        assert "not found" in result.content.lower()

    @pytest.mark.asyncio
    async def test_read_wrong_extension(self, tmp_dir: Path):
        wrong = tmp_dir / "notes.txt"
        wrong.write_text("hello")
        result = await _read_docx(str(wrong))
        assert result.is_error
        assert ".docx" in result.content


class TestXlsxRoundTrip:
    @pytest.mark.asyncio
    async def test_write_then_read(self, tmp_dir: Path):
        path = tmp_dir / "data.xlsx"
        rows = [["Name", "Score"], ["Alice", 95], ["Bob", 88]]
        result = await _write_xlsx(str(path), rows)
        assert not result.is_error
        assert path.exists()

        read_result = await _read_xlsx(str(path))
        assert not read_result.is_error
        assert "Alice" in read_result.content
        assert "95" in read_result.content
        assert "Bob" in read_result.content

    @pytest.mark.asyncio
    async def test_read_missing_file(self, tmp_dir: Path):
        result = await _read_xlsx(str(tmp_dir / "nope.xlsx"))
        assert result.is_error


class TestPptxRead:
    @pytest.mark.asyncio
    async def test_read_generated_pptx(self, tmp_dir: Path):
        from pptx import Presentation

        prs = Presentation()
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Quarterly Results"
        body = slide.placeholders[1]
        body.text_frame.text = "Revenue is up 12% year over year."

        path = tmp_dir / "deck.pptx"
        prs.save(str(path))

        result = await _read_pptx(str(path))
        assert not result.is_error
        assert "Quarterly Results" in result.content
        assert "Revenue is up 12%" in result.content

    @pytest.mark.asyncio
    async def test_read_missing_file(self, tmp_dir: Path):
        result = await _read_pptx(str(tmp_dir / "nope.pptx"))
        assert result.is_error


class TestFileContextRouting:
    @pytest.mark.asyncio
    async def test_build_file_context_routes_docx_through_parser(self, tmp_dir: Path):
        from elidia.cli.main import _build_file_context

        from docx import Document
        doc = Document()
        doc.add_paragraph("Extracted via ingest parser.")
        path = tmp_dir / "sample.docx"
        doc.save(str(path))

        ctx = await _build_file_context((str(path),))
        assert "Extracted via ingest parser." in ctx
        assert "PK" not in ctx  # not the raw zip bytes leaking through as garbled text
