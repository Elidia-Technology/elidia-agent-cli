import hashlib
import logging
import mimetypes
from pathlib import Path

from elidia.rag.engine import RagEngine

logger = logging.getLogger(__name__)

CODE_EXTENSIONS = {
    ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".go", ".rs", ".c", ".cpp", ".h",
    ".rb", ".php", ".swift", ".kt", ".scala", ".sh", ".bash", ".zsh",
    ".sql", ".yaml", ".yml", ".toml", ".json", ".xml", ".html", ".css", ".scss",
}

TEXT_EXTENSIONS = {
    ".txt", ".md", ".rst", ".org", ".tex", ".log", ".csv", ".tsv",
    ".cfg", ".ini", ".conf", ".env.example", ".gitignore", ".dockerignore",
    ".makefile", ".dockerfile",
}

SKIP_DIRS = {
    ".git", ".svn", ".hg", "__pycache__", "node_modules", ".venv", "venv",
    ".tox", ".mypy_cache", ".pytest_cache", ".ruff_cache", "dist", "build",
    ".next", ".nuxt", ".cache", "coverage", ".idea", ".vscode",
}

MAX_FILE_SIZE = 1024 * 1024  # 1MB


class FileIngestPipeline:
    """Ingests files from disk into the RAG engine."""

    def __init__(self, rag: RagEngine) -> None:
        logger.debug("Entered into FileIngestPipeline.__init__")
        self._rag = rag

    async def ingest_file(
        self,
        path: Path,
        project_path: str = "",
        chunk_size: int = 512,
    ) -> list[str]:
        logger.debug(f"Entered into ingest_file: path={path}")
        path = path.resolve()

        if not path.exists():
            logger.warning(f"File not found: {path}")
            return []
        if not path.is_file():
            logger.warning(f"Not a file: {path}")
            return []
        if path.stat().st_size > MAX_FILE_SIZE:
            logger.warning(f"File too large ({path.stat().st_size} bytes): {path}")
            return []

        content_type = self._detect_content_type(path)
        if content_type is None:
            logger.info(f"Skipping unsupported file type: {path}")
            return []

        text = await self._extract_text(path, content_type)
        if not text or not text.strip():
            return []

        file_hash = hashlib.sha256(text.encode()).hexdigest()[:16]

        return await self._rag.ingest(
            text=text,
            source=str(path),
            content_type=content_type,
            project_path=project_path,
            file_hash=file_hash,
            chunk_size=chunk_size,
        )

    async def _extract_text(self, path: Path, content_type: str) -> str:
        """Extract text from a file, dispatching to the appropriate parser."""
        logger.debug(f"Entered into _extract_text: path={path}, type={content_type}")
        try:
            if content_type == "pdf":
                return _parse_pdf(path)
            elif content_type == "docx":
                return _parse_docx(path)
            elif content_type == "xlsx":
                return _parse_xlsx(path)
            elif content_type == "pptx":
                return _parse_pptx(path)
            elif content_type == "html":
                return _parse_html(path)
            else:
                return path.read_text(encoding="utf-8", errors="replace")
        except UnicodeDecodeError:
            return path.read_text(encoding="latin-1", errors="replace")
        except Exception as e:
            logger.warning(f"Could not extract text from {path}: {e}")
            return ""

    async def ingest_directory(
        self,
        directory: Path,
        project_path: str = "",
        chunk_size: int = 512,
        recursive: bool = True,
    ) -> dict[str, int]:
        logger.debug(f"Entered into ingest_directory: dir={directory}, recursive={recursive}")
        directory = directory.resolve()

        if not directory.exists() or not directory.is_dir():
            return {"files": 0, "chunks": 0, "skipped": 0}

        files_ingested = 0
        total_chunks = 0
        skipped = 0

        if recursive:
            file_iter = self._walk_files(directory)
        else:
            file_iter = (f for f in sorted(directory.iterdir()) if f.is_file())

        for file_path in file_iter:
            if self._should_skip(file_path):
                skipped += 1
                continue

            ids = await self.ingest_file(
                file_path,
                project_path=project_path or str(directory),
                chunk_size=chunk_size,
            )
            if ids:
                files_ingested += 1
                total_chunks += len(ids)

        logger.info(f"Ingested {files_ingested} files, {total_chunks} chunks, skipped {skipped}")
        return {"files": files_ingested, "chunks": total_chunks, "skipped": skipped}

    def _walk_files(self, directory: Path):
        for item in sorted(directory.iterdir()):
            if item.is_dir():
                if item.name in SKIP_DIRS:
                    continue
                yield from self._walk_files(item)
            elif item.is_file():
                yield item

    def _should_skip(self, path: Path) -> bool:
        if any(part in SKIP_DIRS for part in path.parts):
            return True
        if path.name.startswith(".") and path.suffix not in {".env.example"}:
            return True
        if path.stat().st_size > MAX_FILE_SIZE:
            return True
        return self._detect_content_type(path) is None

    def _detect_content_type(self, path: Path) -> str | None:
        suffix = path.suffix.lower()
        name = path.name.lower()

        if suffix in CODE_EXTENSIONS or name in {"makefile", "dockerfile", "rakefile", "gemfile"}:
            return "code"
        if suffix in {".md", ".rst"}:
            return "markdown"
        if suffix in TEXT_EXTENSIONS:
            return "text"

        # Binary / rich document formats
        if suffix == ".pdf":
            return "pdf"
        if suffix == ".docx":
            return "docx"
        if suffix == ".xlsx":
            return "xlsx"
        if suffix == ".pptx":
            return "pptx"
        if suffix in {".html", ".htm"}:
            return "html"

        mime, _ = mimetypes.guess_type(str(path))
        if mime and mime.startswith("text/"):
            return "text"

        return None


# --- Binary format parsers ---

def _parse_pdf(path: Path) -> str:
    """Extract text from PDF using pymupdf (fitz)."""
    logger.debug(f"Entered into _parse_pdf: path={path}")
    try:
        import fitz
        doc = fitz.open(str(path))
        parts: list[str] = []
        for page in doc:
            text = page.get_text()
            if text.strip():
                parts.append(text.strip())
        doc.close()
        return "\n\n".join(parts)
    except ImportError:
        logger.warning("pymupdf not installed — cannot parse PDF. pip install pymupdf")
        return ""
    except Exception as e:
        logger.warning(f"PDF parse failed for {path}: {e}")
        return ""


def _parse_docx(path: Path) -> str:
    """Extract text from DOCX using python-docx."""
    logger.debug(f"Entered into _parse_docx: path={path}")
    try:
        from docx import Document
        doc = Document(str(path))
        parts: list[str] = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text.strip())
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells)
                if row_text.strip():
                    parts.append(row_text.strip())
        return "\n\n".join(parts)
    except ImportError:
        logger.warning("python-docx not installed — cannot parse DOCX. pip install python-docx")
        return ""
    except Exception as e:
        logger.warning(f"DOCX parse failed for {path}: {e}")
        return ""


def _parse_xlsx(path: Path) -> str:
    """Extract text from XLSX using openpyxl."""
    logger.debug(f"Entered into _parse_xlsx: path={path}")
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        parts: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"[Sheet: {sheet_name}]")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    parts.append(row_text.strip())
        wb.close()
        return "\n".join(parts)
    except ImportError:
        logger.warning("openpyxl not installed — cannot parse XLSX. pip install openpyxl")
        return ""
    except Exception as e:
        logger.warning(f"XLSX parse failed for {path}: {e}")
        return ""


def _parse_pptx(path: Path) -> str:
    """Extract text from PPTX using python-pptx."""
    logger.debug(f"Entered into _parse_pptx: path={path}")
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts: list[str] = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_parts.append(para.text.strip())
            if slide_parts:
                parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_parts))
        return "\n\n".join(parts)
    except ImportError:
        logger.warning("python-pptx not installed — cannot parse PPTX. pip install python-pptx")
        return ""
    except Exception as e:
        logger.warning(f"PPTX parse failed for {path}: {e}")
        return ""


def _parse_html(path: Path) -> str:
    """Extract visible text from HTML using built-in html.parser."""
    logger.debug(f"Entered into _parse_html: path={path}")
    try:
        from html.parser import HTMLParser

        class TextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self._parts: list[str] = []
                self._skip_tags = {"script", "style", "noscript", "meta", "link"}

            def handle_data(self, data: str) -> None:
                text = data.strip()
                if text:
                    self._parts.append(text)

        extractor = TextExtractor()
        extractor.feed(path.read_text(encoding="utf-8", errors="replace"))
        return "\n".join(extractor._parts)
    except Exception as e:
        logger.warning(f"HTML parse failed for {path}: {e}")
        return ""
